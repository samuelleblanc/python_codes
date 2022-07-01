#!/usr/bin/env python
# coding: utf-8

# In[ ]:


def __init__():
    """
       Collection of codes to run some typical plotting utilities
       
           - circles : plot circles with radius in data scale
           
        details are in the info of each module
    """
    pass


# In[ ]:


def circles(x, y, s, c='b', ax=None, vmin=None, vmax=None, **kwargs):
    """
    Make a scatter of circles plot of x vs y, where x and y are sequence 
    like objects of the same lengths. The size of circles are in data scale.

    Parameters
    ----------
    x,y : scalar or array_like, shape (n, )
        Input data
    s : scalar or array_like, shape (n, ) 
        Radius of circle in data scale (ie. in data unit)
    c : color or sequence of color, optional, default : 'b'
        `c` can be a single color format string, or a sequence of color
        specifications of length `N`, or a sequence of `N` numbers to be
        mapped to colors using the `cmap` and `norm` specified via kwargs.
        Note that `c` should not be a single numeric RGB or
        RGBA sequence because that is indistinguishable from an array of
        values to be colormapped.  `c` can be a 2-D array in which the
        rows are RGB or RGBA, however.
    ax : Axes object, optional, default: None
        Parent axes of the plot. It uses gca() if not specified.
    vmin, vmax : scalar, optional, default: None
        `vmin` and `vmax` are used in conjunction with `norm` to normalize
        luminance data.  If either are `None`, the min and max of the
        color array is used.  (Note if you pass a `norm` instance, your
        settings for `vmin` and `vmax` will be ignored.)

    Returns
    -------
    paths : `~matplotlib.collections.PathCollection`

    Other parameters
    ----------------
    kwargs : `~matplotlib.collections.Collection` properties
        eg. alpha, edgecolors, facecolors, linewidths, linestyles, norm, cmap

    Examples
    --------
    a = np.arange(11)
    circles(a, a, a*0.2, c=a, alpha=0.5, edgecolor='none')

    License
    --------
    This code is under [The BSD 3-Clause License]
    (http://opensource.org/licenses/BSD-3-Clause)
    """
    from matplotlib.patches import Circle
    from matplotlib.collections import PatchCollection
    import pylab as plt
    #import matplotlib.colors as colors

    if ax is None:
        ax = plt.gca()    

    if isinstance(c,str):
        color = c     # ie. use colors.colorConverter.to_rgba_array(c)
    else:
        color = None  # use cmap, norm after collection is created
    kwargs.update(color=color)

    if isinstance(x, (int, float)):
        patches = [Circle((x, y), s),]
    elif isinstance(s, (int, float)):
        patches = [Circle((x_,y_), s) for x_,y_ in zip(x,y)]
    else:
        patches = [Circle((x_,y_), s_) for x_,y_,s_ in zip(x,y,s)]
    collection = PatchCollection(patches, **kwargs)

    if color is None:
        collection.set_array(np.asarray(c))
        if vmin is not None or vmax is not None:
            collection.set_clim(vmin, vmax)

    ax.add_collection(collection)
    return collection


# In[5]:


def plot_color_maps(reverse=False):
    """
    Simple plotting function to run through and plot each color map
    Help for choosing which colormap to use
    """
    import pylab as plt
    from numpy import outer
    plt.rc('text', usetex=False)
    a=outer(plt.ones(10,),plt.arange(0,1,0.01))
    plt.figure(figsize=(5,15))
    plt.subplots_adjust(top=0.8,bottom=0.08,left=0.03,right=0.99)
    if reverse:
        maps=[m for m in plt.cm.datad]
        rr = 2
    else:
        maps=[m for m in plt.cm.datad if not m.endswith("_r")]
        rr = 1
    maps.sort()
    l=len(maps)+1
    title_dict = {'fontsize': 10,
                  'verticalalignment': 'center',
                  'horizontalalignment': 'left'}
    for i, m in enumerate(maps):
        plt.subplot(l,rr,i+1)
        plt.axis("off")
        plt.imshow(a,aspect='auto',cmap=plt.get_cmap(m),origin="lower")
        plt.text(1.01,0.5,m,fontdict=title_dict,transform=plt.gca().transAxes)


# In[1]:


def plot_vert_hist(fig,ax1,y,pos,ylim,color='grey',label=None,legend=False,onlyhist=True,loc=2,bins=30,alpha=0.5):
    """
    Purpose:
        function to plot a 'bean' like vertical histogram
    
    Input:
        fig: figure object for reference and plotting
        ax1: axis object to plot on
        y: values to be plotted in histogram fashion
        pos: position along axis to create the bean plot
        ylim: range of plotted axis [min,max]
        color: (default to grey) color of the plotted histogram
        label: (default to none) if included, the label for the color squares of this histogram
        legend: (default to False) If True adds a legend on the location 
        onlyhist: (default to True) only plots the histogram bean, not the mean and median lines
        loc: (default to 2) location of the legend
        bins: (default to 30) number of bins to plot
        alpha: (defautl to 0.5) value of transparency (0 to 1)
    
    Output:
        only the histogram plot on top of an existing plot
    
    Dependencies:
        - numpy
        - Sp_parameters
        - plotting_utils : this file
    
    Required files:
        None
    
    Modification History:
        Written: Samuel LeBlanc, NASA Ames
        Modified: Samuel LeBlanc, NASA Ames in Santa Cruz, 2015-12-09
                  - added comments
                  - added the bins keyword to be used
        Modified: Samuel LeBlanc, NASA Ames in Santa Cruz, 2015-12-12
                  - added alpha keyword
    """
    import Sp_parameters as Sp
    import numpy as np
    from plotting_utils import data2figpoints
    (ymask,iy) = Sp.nanmasked(y)
    ax = fig.add_axes(data2figpoints(pos,0.4,fig=fig,ax1=ax1),frameon=False,ylim=ylim)
    ax.tick_params(axis='both', which='both', labelleft='off', labelright='off',bottom='off',top='off',
               labelbottom='off',labeltop='off',right='off',left='off')
    ax.hist(ymask,orientation='horizontal',normed=True,color=color,edgecolor='None',bins=bins,alpha=alpha,label=label,range=ylim)
    if onlyhist:
        label_mean = None
        label_median = None
    else:
        label_mean = 'Mean'
        label_median = 'Median'
    ax.axhline(np.mean(ymask),color='red',linewidth=2,label=label_mean)
    ax.axhline(np.median(ymask),color='k',linewidth=2,linestyle='--',label=label_median)
    if legend:
        ax.legend(frameon=False,loc=loc)
    ax = fig.add_axes(data2figpoints(pos+0.01,-0.4,fig=fig,ax1=ax1),frameon=False,ylim=ylim)
    ax.tick_params(axis='both', which='both', labelleft='off', labelright='off',bottom='off',top='off',
                   labelbottom='off',labeltop='off',right='off',left='off')
    ax.hist(ymask,orientation='horizontal',normed=True,color=color,edgecolor='None',bins=bins,alpha=alpha,range=ylim)
    ax.axhline(np.mean(ymask),color='red',linewidth=2)
    ax.axhline(np.median(ymask),color='k',linewidth=2,linestyle='--')


# In[ ]:


def data2figpoints(x,dx,fig,ax1):
    "function to tranform data locations to relative figure coordinates (in fractions of total figure"
    flen = fig.transFigure.transform([1,1])
    bot = ax1.transAxes.transform([0,0])/flen
    top = ax1.transAxes.transform([1,1])/flen
    
    start = ax1.transData.transform([x,0])/flen
    end = ax1.transData.transform([x+dx,0])/flen
    left = start[0]
    bottom = bot[1]
    width = end[0]-start[0]
    height = top[1]-bot[1] 
    return left,bottom,width,height


# In[3]:


def plot_lin(x,y,x_err=[None],y_err=[None],color='b',labels=True,ci=0.95,
             shaded_ci=True,use_method='linfit',ax=None,lblfmt='2.2f',label_prefix='',*args,**kwargs):
    """
    function to plot on top of previous a linear fit line, 
    with the line equation in legend.
    Input:
       x: independent
       y: dependent
       x_err: uncertainty in x (default None)
       y_err: uncertainty in y (default None)
       color: color of the plot (default blue)
       labels: if include label in legend of linear equation values (default True)
       lblfmt: format of labels (default is 2.2f) 
       ci: Confidence interval (in percent) (default 95)
       shaded_ci: plot the shaded confidence interval (default True)
       use_method: Define which method to use for linear regression
                   options:
                   'linfit' (default) Use the linfit method from linfit module, when set, x_err and y_err are ignored
                   'odr' use the scipy ODR method to calculate the linear regression, with x_err and y_err abilities
                   'statsmodels' use the statsmodels method, Weighted least squares, with weighing of 1/y_err, x_err ignored
                   'york' Use the bivariate_fit defined by York et al. (2004)
       ax: variable containing the axis to which to plot onto.
       label_prefix
       any other input for matplotlib plot function can be passed via args or kwargs
       
    Output:
        p coefficients (intercept, slope)
        perr values (error in intercept, error in slope)
    """
    import matplotlib.pyplot as plt
    import numpy as np
    from Sp_parameters import doublenanmask, nanmasked
    from plotting_utils import confidence_envelope, lin
    if not ax:
        ax = plt.gca()
    xn,yn,mask = doublenanmask(x,y,return_mask=True)
    if label_prefix: labels=True
    if use_method=='odr':
        from scipy import odr
        model = odr.Model(lin)
        if any(x_err):
            if any(y_err):
                dat = odr.RealData(xn,yn,sx=x_err[mask],sy=y_err[mask])
            else:
                dat = odr.RealData(xn,yn,sx=x_err[mask])
        else:
            if any(y_err):
                dat = odr.RealData(xn,yn,sy=y_err[mask]) 
            else:
                dat = odr.RealData(xn,yn)
        try:
            from linfit import linfit
            c,cm = linfit(xn,yn)
            p = np.array([c[1],c[0]])
        except:
            p = [1.0,0.5]
        outa = odr.ODR(dat,model,beta0=p).run()
        print(outa.cov_beta)
        perr = np.sqrt(np.diag(outa.cov_beta))
        p = outa.beta
    elif use_method=='linfit':
        from linfit import linfit
        c,cm = linfit(xn,yn)
        p = np.array([c[1],c[0]])
        cerr = np.sqrt(np.diag(cm))
        perr = np.array([cerr[1],cerr[0]]) 
    elif use_method=='statsmodels':
        import statsmodels.api as sm
        Xn = sm.add_constant(xn)
        if any(y_err):
            results = sm.WLS(yn,Xn,weights=1/y_err[mask]).fit()
        else:
            results = sm.OLS(yn,Xn).fit()
        p = results.params
        perr = results.bse
    elif use_method=='york':
        from plotting_utils import bivariate_fit
        try:
            from linfit import linfit
            c,cm = linfit(xn,yn)
            p = np.array([c[1],c[0]])
        except:
            p = [1.0,0.5]
        if not any(x_err):
            raise('x_err must be set for york fit')
        if not any(y_err):
            raise('y_err must be set for york fit')
        ri = np.corrcoef(x_err[mask],y_err[mask])[0,1]**2
        a_bivar, b_bivar, S, cov = bivariate_fit(xn,yn,x_err[mask],y_err[mask],b0=p[1],ri=ri)
        p = [a_bivar,b_bivar]
        cerr = np.sqrt(np.diag(cov))
        perr = np.array([cerr[1],cerr[0]]) 
    else:
        print('Method: %s is not a valid choice' % use_method)
        return
    xx = np.linspace(xn.min()-np.abs(xn.min()*0.1),xn.max()+np.abs(xn.max()*0.1))
    if labels:
        ax.plot(xx,lin(p,xx),color=color,
                label='{label_prefix}y=({:{fmt}}$\pm${:{fmt}})+\n({:{fmt}}$\pm${:{fmt}})x'.format(
                        p[0],perr[0],p[1],perr[1],fmt=lblfmt,label_prefix=label_prefix),*args,**kwargs)
    else:
        ax.plot(xx,lin(p,xx),color=color,*args,**kwargs)
    if shaded_ci:
        y_up,y_down = confidence_envelope(xx, p, perr, ci=ci)
        ax.fill_between(xx,y_down,y_up,color=color,alpha=0.1)
    return p,perr


# In[ ]:


def lin(p,x):
    """ 
    Simple function that returns a linear expression:
    y = p[0] + p[1]*x
    """
    return p[0]+p[1]*x


# In[ ]:


def confidence_envelope(xn,p,p_err,ci=95,size=1000):
    """
    Simple function to model the confidence enveloppe of a linear function
    Returns y_up and y_down: y values for the upper bounds and lower bound
       of the confidence interval defined by ci [in percent]. Returns the y values at each xn points
    Inputs:
    p: p[0] is the intercept, p[1] is the slope
    p_err: uncertainty in each p value
    size: number of points to use in montecarlo simulation of confidence bounds (default 1000)
    """
    import numpy as np
    from plotting_utils import lin
    from scipy import stats
    if len(xn)<=1:
        print('** Problem with input xn **')
        return None,None
    p0s = np.random.normal(loc=p[0],scale=p_err[0],size=size)
    p1s = np.random.normal(loc=p[1],scale=p_err[1],size=size)
    ys = np.zeros((size,len(xn)))
    for i in range(size):
        ys[i,:] = lin([p0s[i],p1s[i]],xn)
    y_up, y_down = np.zeros((2,len(xn)))
    for j in range(len(xn)):
        y_up[j] = stats.scoreatpercentile(ys[:,j],ci)
        y_down[j] = stats.scoreatpercentile(ys[:,j],100-ci)
    return y_up, y_down


# In[ ]:


def plotmatfig(filename,fignr=None):
    """
    Plot a figure from a matlab .fig file
    
    Taken from the http://stackoverflow.com/questions/8172931/data-from-a-matlab-fig-file-using-python
    user response: johnml1135
    
    """
    
    from scipy.io import loadmat
    import numpy as np
    import matplotlib.pyplot as plt
    d = loadmat(filename,squeeze_me=True, struct_as_record=False)
    matfig = d['hgS_070000']
    childs = matfig.children
    ax1 = [c for c in childs if c.type == 'axes']
    multi,lasta = False,-999
    if(len(ax1) > 0):
        for i,a in enumerate(ax1):
            if type(a.children) is np.ndarray:
                if ax1==lasta:
                    multi = True
                    ax2 = [ax1,a]
                elif multi:
                    ax2.append(a)
                else:
                    ax1 = a
                    ax2 = [ax1]
                lasta = a
    legs = [c for c in childs if c.type == 'scribe.legend']
    if(len(legs) > 0):
        legs = legs[0]
    else:
        legs=0
    pos = matfig.properties.Position
    size = np.array([pos[2]-pos[0],pos[3]-pos[1]])/96
    XX,YY = [],[]
    for ax1 in ax2:
        plt.figure(fignr,figsize=size)
        plt.clf()
        #plt.hold(True)
        counter = 0    
        for line in ax1.children:
            if line.type == 'graph2d.lineseries':
                if hasattr(line.properties,'Marker'):
                    mark = "%s" % line.properties.Marker
                    if(mark != "none"):
                        mark = mark[0]
                else:
                    mark = '.'
                if hasattr(line.properties,'LineStyle'):
                    linestyle = "%s" % line.properties.LineStyle
                else:
                    linestyle = '-'
                if hasattr(line.properties,'Color'):
                    r,g,b =  line.properties.Color
                else:
                    r = 0
                    g = 0
                    b = 1
                if hasattr(line.properties,'MarkerSize'):
                    marker_size = line.properties.MarkerSize
                else:
                    marker_size = -1                
                x = line.properties.XData
                y = line.properties.YData
                XX.append(x)
                YY.append(y)
                if(mark == "none"):
                    plt.plot(x,y,linestyle=linestyle,color=[r,g,b])
                elif(marker_size==-1):
                    plt.plot(x,y,marker=mark,linestyle=linestyle,color=[r,g,b])
                else:
                    plt.plot(x,y,marker=mark,linestyle=linestyle,color=[r,g,b],ms=marker_size)
            elif line.type == 'text':
                if counter == 1:
                    try:
                        plt.xlabel("$%s$" % line.properties.String,fontsize =16)
                    except:
                        pass
                elif counter == 2:
                    try:
                        plt.ylabel("$%s$" % line.properties.String,fontsize = 16)
                    except:
                        pass
                elif counter == 4:
                    try:
                        plt.title("$%s$" % line.properties.String,fontsize = 16)
                    except:
                        pass
                counter += 1   
        plt.grid(ax1.properties.__dict__.get('XGrid'))

        if(hasattr(ax1.properties,'XTick')):
            if(hasattr(ax1.properties,'XTickLabelRotation')):
                plt.xticks(ax1.properties.XTick,ax1.properties.XTickLabel,rotation=ax1.properties.XTickLabelRotation)
            elif(hasattr(ax1.properties,'XTickLabel')):
                plt.xticks(ax1.properties.XTick,ax1.properties.XTickLabel)
        if(hasattr(ax1.properties,'YTick')):
            if(hasattr(ax1.properties,'YTickLabelRotation')):
                plt.yticks(ax1.properties.YTick,ax1.properties.YTickLabel,rotation=ax1.properties.YTickLabelRotation)
            elif(hasattr(ax1.properties,'YTickLabel')):
                plt.yticks(ax1.properties.YTick,ax1.properties.YTickLabel)
        if(hasattr(ax1.properties,'XLim')):
            plt.xlim(ax1.properties.XLim)
        if(hasattr(ax1.properties,'YLim')):
            plt.ylim(ax1.properties.YLim)
        if legs:        
            leg_entries = tuple(['$' + l + '$' for l in legs.properties.String])
            py_locs = ['upper center','lower center','right','left','upper right','upper left','lower right','lower left','best','best']
            MAT_locs=['North','South','East','West','NorthEast', 'NorthWest', 'SouthEast', 'SouthWest','Best','none']
            Mat2py = dict(list(zip(MAT_locs,py_locs)))
            location = legs.properties.Location
            plt.legend(leg_entries,loc=Mat2py[location])
        #plt.hold(False)
        plt.show()
    return XX,YY


# In[ ]:


def make_pptx(filepath,filename,title='',glob_pattern='*',wide=False):
    """
    Purpose:
        function to make a powerpoint presentation with all figures within a single folder, following a glob pattern
    
    Input:
        filepath: full path to where to find the images or other files
        filename: filename of the pptx file to be created
        title: the title of the presentation
        glob_pattern: (defaults to '*') the pattern to be used to discern which file to include
        wide: (default to False) if set to True, then outputs a pptx in the widescreen format (16:9)
              otherwise, the standard 4:3 format
    
    Output:
        pptx file under the path filepath/filename.pptx
    
    Dependencies:
        - pptx
        - glob
        - scipy
        - datetime
    
    Required files:
        None
    
    Modification History:
        Written: Samuel LeBlanc, NASA Ames, CA, 2016-10-20
                ported from pptximage.py at https://gist.github.com/glass5er/748cda36befe17fd1cb0, user glass5er
    """
    import pptx
    import pptx.util
    import glob
    import scipy.misc
    from datetime import date

    OUTPUT_TAG = title

    prs = pptx.Presentation()

    # default slide width
    prs.slide_width = 9144000
    
    if wide:
        # slide height @ 16:9
        prs.slide_height = 5143500
    else:
        # slide height @ 4:3
        prs.slide_height = 6858000
    
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # blank slide
    #slide = prs.slides.add_slide(prs.slide_layouts[6])

    # set title
    title = slide.shapes.title
    title.text = OUTPUT_TAG
    subtitle = slide.placeholders[1]
    subtitle.text = "Generated on {:%Y-%m-%d}".format(date.today())

    pic_left  = int(prs.slide_width * 0.05)
    pic_top   = int(prs.slide_height * 0.1)
    pic_width = int(prs.slide_width * 0.9)

    for g in glob.glob(filepath+glob_pattern):
        pic_left  = int(prs.slide_width * 0.05)
        pic_width = int(prs.slide_width * 0.9)
        print(g)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
        p = tb.textframe.add_paragraph()
        p.text = g.split('\\')[-1]
        p.font.size = pptx.util.Pt(14)

        try:
            img = scipy.misc.imread(g)
            pic_height = int(pic_width * img.shape[0] / img.shape[1])
            if pic_height>prs.slide_height:
                h,w = pic_height,pic_width
                pic_height = int(prs.slide_height * 0.9)
                pic_width = int(pic_height * w/h)
                pic_left = int((prs.slide_width-pic_width)/2 + prs.slide_width * 0.05)
                #import pdb; pdb.set_trace()
        except:
            print('Error on picture: {} using default size values'.format(g)) 
        #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
        pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
    print('Saving to: {}{}.pptx'.format(filepath,filename))
    prs.save(filepath+'%s.pptx' % filename)


# In[1]:


def color_box(bp, color):
    'Coloring of all the elements of a box plot'
    import matplotlib.pyplot as plt
    
    # Define the elements to color. You can also add medians, fliers and means
    elements = ['boxes','caps','whiskers','medians','means','fliers']

    if type(color) is not list: color = [color]
    # Iterate over each of the elements changing the color
    if len(color) == len(bp[elements[0]]):
        colors = color
    elif len(color)==1:
        colors = []
        [colors.extend(color) for i in range(len(bp[elements[0]]))]
        
    for elem in elements:
        if len(bp[elem]) > len(colors):
            [plt.setp(bp[elem][idx], color=colors[idx/2]) for idx in range(len(bp[elem]))]
        else:
            [plt.setp(bp[elem][idx], color=colors[idx]) for idx in range(len(bp[elem]))]
    return


# In[ ]:


def subset_bins(vals,val_lim,lims):
    'create the subsetted bins of values'
    bins = []
    for i,c in enumerate(lims[0:-1]):
        val_fl = (val_lim>=c)&(val_lim<lims[i+1])
        bins.append(vals[val_fl])
    return bins


# In[1]:


def make_boxplot(vals,val_lim,lims,pos,color='green',label=None,y=0,alpha=1.0, ax=None,vert=True,fliers_off=False,
                 tick_labels=True,return_bp=False,mean_marker='s',**kwargs):
    """Compile the functions to make a box plot
    
    vals: values to box
    val_lim: values to use as basis for binning
    lims: limits of the bins to use
    pos: center position of the limites
    y:?
    vert: (default True) if True, return vertical boxes, false for horizontal boxes
    fliers_off: (default False) if True, turns off the plotting of the outliers
    return_bp: (default False) returns the boxplot links if True
    mean_marker: (default 's') the marker for the mean point 
    """

    import matplotlib.pyplot as plt
    from plotting_utils import subset_bins, color_box
    
    if not ax:
        ax = plt.gca()
        
    if vert:
        ti = ax.get_xticks()
        tl = ax.get_xticklabels()
    else:
        ti = ax.get_yticks()
        tl = ax.get_yticklabels()
        
    bins = subset_bins(vals,val_lim,lims)
    
    bo = ax.boxplot(bins,y,'.',showmeans=True,positions=pos,vert=vert,**kwargs)
    color_box(bo,color)
    for n in list(bo.keys()):
        nul = [plt.setp(bo[n][idx],alpha=alpha)for idx in range(len(bo[n]))]
    if fliers_off:
        u = [plt.setp(bo['fliers'][idx],alpha=0.00)for idx in range(len(bo['fliers']))]
    else:
        u = [plt.setp(bo['fliers'][idx],alpha=0.04)for idx in range(len(bo['fliers']))]
    v = [plt.setp(bo['means'][idx],alpha=0.05)for idx in range(len(bo['means']))]
    if vert:
        mean = [a.get_ydata()[0] for a in bo['means']]
        ax.plot(pos, mean,mean_marker+'-',zorder=100,color=color,label=label,lw=2.5,alpha=alpha)
    else:
        mean = [a.get_xdata()[0] for a in bo['means']]
        ax.plot( mean,pos,mean_marker+'-',zorder=100,color=color,label=label,lw=2.5,alpha=alpha)
        
    
    #plt.gca().xaxis.set_major_locator(AutoLocator())
    #plt.gca().xaxis.set_major_locator(AutoLocator)
    
    if vert:
        ti1 = ax.set_xticks(ti)
        if tick_labels: 
            tl1 = ax.set_xticklabels([t for t in tl])
        else: 
            ax.set_xticks([])
            ax.set_xticklabels([])
            
    else:
        ti1 = ax.set_yticks(ti)
        if tick_labels: 
            tl1 = ax.set_yticklabels([t for t in tl])
        else: 
            ax.set_yticks([])
            ax.set_yticklabels([])
    if return_bp:
        return mean, bo
    else:
        return mean


# In[ ]:


def prelim(ax=None):
    'Stamp prelim in center of the plot'
    import matplotlib.pyplot as plt
    if not ax:
        ax = plt.gca()
    
    ax.text(0.5, 0.5, 'Preliminary',
        verticalalignment='bottom', horizontalalignment='center',
        transform=ax.transAxes,
        color='k', fontsize=18,zorder=1,alpha=0.3)


# In[ ]:


def sub_note(note,ax=None,out=False,dx=0.0,dy=0.0,fontsize=18):
    'Stamp note in top right of the plot, adjust with dx, dy, if out set to true, is put outside the plot'
    import matplotlib.pyplot as plt
    if not ax:
        ax = plt.gca()
    
    if out: 
        yup = 1.02
        val = 'bottom'
    else: 
        yup = 0.98
        val = 'top'
    ax.text(0.01+dx, yup+dy, ' '+note,
        verticalalignment=val, horizontalalignment='left',
        transform=ax.transAxes,
        color='k', fontsize=fontsize,zorder=1,alpha=0.7)


# In[1]:


def set_box_whisker_color(cl,bp,binned_ndays,color_not_start_at_zero=False,
                          mean_color='darkgreen',whisker_color='pink',median_color='gold',face_alpha=1.0):
    'To change the color (cl=colormap) of box and whisker plots (bp=box_whisker plot artists) to denote the number of samples (binned_ndays=number of samples), if colors dont start at zero, set color_not_start_at_zero to True' 
    import numpy as np
    bndm = np.nanmax(binned_ndays)*1.0
    if color_not_start_at_zero: 
        minb = np.nanmin(binned_ndays)*1.0
        bndm = np.nanmax(binned_ndays)*1.0 - minb
    for j,b in enumerate(bp['boxes']):
        if color_not_start_at_zero:
            b.set_facecolor(cl((binned_ndays[j]*1.0-minb)/bndm))
            b.set_edgecolor(cl((binned_ndays[j]*1.0-minb)/bndm))
        else:
            b.set_facecolor(cl(binned_ndays[j]*1.0/bndm))
            b.set_edgecolor(cl(binned_ndays[j]*1.0/bndm))
        b.set_alpha(face_alpha)
    for j,b in enumerate(bp['means']):
        b.set_marker('.')
        b.set_color('None')
        b.set_markerfacecolor(mean_color)
        b.set_markeredgecolor(mean_color)
        b.set_alpha(0.6)
    for j,b in enumerate(bp['whiskers']):
        b.set_linestyle('-')
        b.set_color(whisker_color) #gr(binned_ndays[j]*1.0/bndm))
        b.set_alpha(0.7)
    for j,b in enumerate(bp['caps']):
        b.set_alpha(0.7)
        b.set_color(whisker_color)#gr(binned_ndays[j]*1.0/bndm))
    for j,b in enumerate( bp['medians']):
        b.set_linewidth(4)
        b.set_color(median_color)
        b.set_alpha(0.4)
    
    return


# In[1]:


def match_ygrid(ax1,ax2,ticks):
    'function to match the grid ticks to a dual y axis plot, matching limits of ax2 such that the ticks line up with ax1 grid'
    y0,y1 = ax1.get_ybound()
    ti = ax1.get_yticks()
    ax2.set_yticks(ticks)
    if ax1.get_yscale() =='log':
        a = (np.log10(ti[1])-np.log10(ti[0]))/(ticks[1]-ticks[0])
        dy = a*ticks[0]-np.log10(ti[0])
        ax2.set_ylim((np.log10(y0)+dy)/a,(np.log10(y1)+dy)/a)
    else:
        a = (ti[1]-ti[0])/(ticks[1]-ticks[0])
        dy = a*ticks[0]-ti[0]
        ax2.set_ylim((y0+dy)/a,(y1+dy)/a)


# In[1]:


def bivariate_fit(xi, yi, dxi, dyi, ri=0.0, b0=1.0, maxIter=1e6):
    """Make a linear bivariate fit to xi, yi data using York et al. (2004).

    This is an implementation of the line fitting algorithm presented in:
    York, D et al., Unified equations for the slope, intercept, and standard
    errors of the best straight line, American Journal of Physics, 2004, 72,
    3, 367-375, doi = 10.1119/1.1632486

    See especially Section III and Table I. The enumerated steps below are
    citations to Section III

    Parameters:
      xi, yi      x and y data points
      dxi, dyi    errors for the data points xi, yi
      ri          correlation coefficient for the weights
      b0          initial guess b
      maxIter     float, maximum allowed number of iterations

    Returns:
      a           y-intercept, y = a + bx
      b           slope
      S           goodness-of-fit estimate
      sigma_a     standard error of a
      sigma_b     standard error of b

    Usage:
    [a, b] = bivariate_fit( xi, yi, dxi, dyi, ri, b0, maxIter)

    """
    import numpy as np
    # (1) Choose an approximate initial value of b
    b = b0

    # (2) Determine the weights wxi, wyi, for each point.
    wxi = 1.0 / dxi**2.0
    wyi = 1.0 / dyi**2.0

    alphai = (wxi * wyi)**0.5
    b_diff = 999.0

    # tolerance for the fit, when b changes by less than tol for two
    # consecutive iterations, fit is considered found
    tol = 1.0e-8

    # iterate until b changes less than tol
    iIter = 1
    while (abs(b_diff) >= tol) & (iIter <= maxIter):

        b_prev = b

        # (3) Use these weights wxi, wyi to evaluate Wi for each point.
        Wi = (wxi * wyi) / (wxi + b**2.0 * wyi - 2.0*b*ri*alphai)

        # (4) Use the observed points (xi ,yi) and Wi to calculate x_bar and
        # y_bar, from which Ui and Vi , and hence betai can be evaluated for
        # each point
        x_bar = np.sum(Wi * xi) / np.sum(Wi)
        y_bar = np.sum(Wi * yi) / np.sum(Wi)

        Ui = xi - x_bar
        Vi = yi - y_bar

        betai = Wi * (Ui / wyi + b*Vi / wxi - (b*Ui + Vi) * ri / alphai)

        # (5) Use Wi, Ui, Vi, and betai to calculate an improved estimate of b
        b = np.sum(Wi * betai * Vi) / np.sum(Wi * betai * Ui)

        # (6) Use the new b and repeat steps (3), (4), and (5) until successive
        # estimates of b agree within some desired tolerance tol
        b_diff = b - b_prev

        iIter += 1

    # (7) From this final value of b, together with the final x_bar and y_bar,
    # calculate a from
    a = y_bar - b * x_bar

    # Goodness of fit
    S = np.sum(Wi * (yi - b*xi - a)**2.0)

    # (8) For each point (xi, yi), calculate the adjusted values xi_adj
    xi_adj = x_bar + betai

    # (9) Use xi_adj, together with Wi, to calculate xi_adj_bar and thence ui
    xi_adj_bar = np.sum(Wi * xi_adj) / np.sum(Wi)
    ui = xi_adj - xi_adj_bar

    # (10) From Wi , xi_adj_bar and ui, calculate sigma_b, and then sigma_a
    # (the standard uncertainties of the fitted parameters)
    sigma_b = np.sqrt(1.0 / np.sum(Wi * ui**2))
    sigma_a = np.sqrt(1.0 / np.sum(Wi) + xi_adj_bar**2 * sigma_b**2)

    # calculate covariance matrix of b and a (York et al., Section II)
    cov = -xi_adj_bar * sigma_b**2
    # [[var(b), cov], [cov, var(a)]]
    cov_matrix = np.array(
        [[sigma_b**2, cov], [cov, sigma_a**2]])

    if iIter <= maxIter:
        return a, b, S, cov_matrix
    else:
        print("bivariate_fit.py exceeded maximum number of iterations, " +
              "maxIter = {:}".format(maxIter))
        return np.nan, np.nan, np.nan, np.nan


# In[5]:


def stats_label(x,y,fmt='2.2f'):
    'To make labels consistently of the relationship between two variables'
    from sklearn.metrics import mean_squared_error, mean_absolute_error
    import scipy.stats as st
    import numpy as np
    
    fl = np.isfinite(x) & np.isfinite(y)
    
    r = st.spearmanr(x,y,nan_policy='omit')
    rmse = mean_squared_error(x[fl],y[fl],squared=True)
    mae = mean_absolute_error(x[fl],y[fl])
    
    return 'R$_{{spearman}}$={:{fmt}}\nRMSE={:{fmt}}\nMAE={:{fmt}}'.format(r.correlation,rmse,mae,fmt=fmt)

